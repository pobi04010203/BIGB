"""주제 설명용 3장 요약 ― 지도교수 대상. 표지·목차·참고문헌 없음.

내용은 ADDENDUM-01 의 재포지셔닝을 따른다:
  방법론은 기존 것이고, 우리는 가정된 감쇠 함수 자리에 실측값을 채운다.
  "최초·novel" 류 표현 금지. 중첩 반전을 우리 발견으로 쓰지 않는다.
  §5.4 의 3단 비교(기하 / 가정 곡선 / 실측 곡선)가 기여를 증명하는 그림이다.

디자인은 hallmark(anti-AI-slop) editorial 규칙을 따른다:
  - 슬라이드마다 구조를 다르게 한다 (같은 골격 반복 = 템플릿 티)
  - 카드·라운드 박스 금지. 헤어라인으로 구분한다
  - 번호를 제목 왼쪽에 붙이는 2단 헤더 금지 (게이트 54). 쪽번호는 하단 folio 로
  - 순백·순흑 금지, 강조색 하나, 좌우 여백 비대칭

폰트 함정:
  1. python-pptx 의 font.name 은 <a:latin> 만 채운다. 한글은 <a:ea> 를 본다.
     ea 가 비면 PowerPoint 에서 테마 기본값으로 나온다. LibreOffice 렌더로는 안 잡힌다.
  2. KoPub 은 굵기마다 패밀리 이름이 다르다. bold 플래그를 쓰면 가짜 볼드가 합성된다.
"""
import csv
import json
import math
import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

import fonts
import theme as T

ROOT = Path(__file__).resolve().parents[1]
FIGS = ROOT / "outputs" / "figures"
OUT = Path(os.environ.get("BRIEF_OUT") or (ROOT / "outputs" / "주제설명_3장.pptx"))

cp = json.loads((ROOT / "outputs" / "curve_params.json").read_text(encoding="utf-8"))
cm = json.loads((ROOT / "outputs" / "comparison.json").read_text(encoding="utf-8"))
sn = json.loads((ROOT / "outputs" / "sensitivity.json").read_text(encoding="utf-8"))
se = json.loads((ROOT / "outputs" / "site_eval.json").read_text(encoding="utf-8"))

PL = cm["placements"]
GEO, ASM, EMP = PL["geometric"], PL["assumed"], PL["empirical"]
D_EA = cm["delta_WDR"]["empirical_minus_assumed"] * 100      # 실측 - 가정
D_EG = cm["delta_WDR"]["empirical_minus_geometric"] * 100    # 실측 - 기하
BASE_PCT = cp["baseline_P"] * 100
DLO, DHI = (v * 100 for v in sn["delta_WDR_range"])
THR_PCT = int(se["threshold"] * 100)

# ── 2번 슬라이드의 숫자는 전부 여기서 유도한다 ─────────────────────────────
# 하드코딩하면 곡선을 다시 피팅했을 때 조용히 틀린 값이 된다 (CLAUDE.md §0.1).
_F, _G = cp["f_rho"], cp["g_theta"]["params"]
_ROWS = list(csv.DictReader((ROOT / "outputs" / "grid_results.csv").open(encoding="utf-8")))


def _measured(rho, theta, occ):
    """격자에서 한 점의 실측 검출률을 기준 조건 대비 비율로 돌려준다."""
    for r in _ROWS:
        if (abs(float(r["rho_px"]) - rho) < 1e-6
                and abs(float(r["theta_deg"]) - theta) < 1e-6
                and abs(float(r["occ_pct_target"]) - occ) < 1e-6):
            return float(r["recall_nohat"]) / cp["baseline_P"] * 100
    raise SystemExit(f"격자에 없는 점: rho={rho} theta={theta} occ={occ}")


def _g(t):
    return ((1 + math.exp(-_G["k"] * _G["x0"]))
            / (1 + math.exp(_G["k"] * (t - _G["x0"]))) * 100)


# 검출률이 기준의 90% 로 떨어지는 머리 크기 -> 4K·HFOV 90° 에서의 거리
RHO_90 = _F["x0"] + math.log(9.0) / _F["k"]
D_90_M = 480.0 / RHO_90        # rho = f_px·H/d 이고 4K·HFOV 90°·머리 0.25m 에서 f_px·H = 480
THETA_HALF = _G["x0"]                                   # 각도 곡선의 반감점
PCT_AT_40M = _measured(12.0, 0.0, 0.0)                  # 40 m = 머리 12px
PCT_AT_45DEG = _g(45.0)
PCT_AT_OCC15 = _measured(48.0, 0.0, 15.0)

# 세 축을 같은 자로 비교한다 ― 각 축이 검출률을 기준의 90% 로 떨어뜨리는 지점.
# 축마다 단위가 달라 "무엇이 더 치명적인가" 를 그냥 견줄 수 없기 때문이다.
THETA_90 = _G["x0"] - math.log(9.0) / _G["k"]
OCC_90 = -math.log(0.9) / cp["h_occ"]["lambda"] * 100

# 현장 규모는 site_eval 에서 읽는다. 격자 계산값(50×30)이 아니라 실제 복셀 수다.
N_VOXEL = len(se["voxels"])
FAIL_PCT_GEO = GEO["fail_voxel_count"] / N_VOXEL * 100
FAIL_PCT_EMP = EMP["fail_voxel_count"] / N_VOXEL * 100
THR_SWEEP = sn["threshold"]

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(T.W), Inches(T.H)
BLANK = prs.slide_layouts[6]


# ── 원시 도구 ──────────────────────────────────────────────────────────────
def _set_face(run, face):
    """latin·ea·cs 를 모두 채운다. ea 가 비면 한글이 테마 폰트로 나온다."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", face)


def text(slide, body, x, y, w, h, size, color=T.BODY, face=T.TEXT,
         align=PP_ALIGN.LEFT, space_after=0, line=None, hang=None):
    """body: 문자열, 또는 [(문자열, {face,size,color}), ...] 런 목록의 문단 리스트."""
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    paras = body if isinstance(body, list) else [body]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line:
            p.line_spacing = line
        if hang:
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(int(Inches(hang))))
            pPr.set("indent", str(-int(Inches(hang))))
        runs = para if isinstance(para, list) else [(para, {})]
        for txt, opt in [(r, {}) if isinstance(r, str) else r for r in runs]:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(opt.get("size", size))
            r.font.color.rgb = opt.get("color", color)
            r.font.bold = opt.get("bold", False)
            _set_face(r, opt.get("face", face))
    return tf


def rule(slide, x, y, w, weight=0.9, color=T.RULE):
    """헤어라인. 카드 테두리 대신 이것만 쓴다."""
    h = weight / 72.0
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return y + h


def vrule(slide, x, y, h, weight=0.9, color=T.RULE):
    w = weight / 72.0
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False


def double_rule(slide, x, y, w):
    """편집 관습의 겹줄. 제목 아래에 쓴다."""
    rule(slide, x, y, w, weight=2.2, color=T.INK)
    rule(slide, x, y + 0.055, w, weight=0.8, color=T.RULE)
    return y + 0.055 + 0.8 / 72.0


def figure(slide, name, x, y, w):
    img = Image.open(FIGS / name)
    slide.shapes.add_picture(str(FIGS / name), Inches(x), Inches(y), width=Inches(w))
    return y + w * img.height / img.width


def page(n):
    """미색 종이 + 하단 folio. 쪽번호는 여기에만 둔다 (제목 옆 번호 금지)."""
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = T.PAPER
    bg.line.fill.background(); bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)
    text(s, f"{n} / 3", T.W - T.MR - 1.2, 6.94, 1.2, 0.3,
         T.FOLIO, T.FAINT, face=T.DISPLAY_LIGHT, align=PP_ALIGN.RIGHT)
    return s


def note(slide, body, y=None):
    text(slide, body, T.ML, y if y is not None else 6.94,
         T.CW - 1.4, 0.34, T.FOLIO, T.FAINT, line=1.3)


# ── 줄 수를 실제로 재서 높이를 구한다 ─────────────────────────────────────
# 좌표를 손으로 맞추면 폰트·문장이 바뀔 때마다 선이 글자를 자른다.
# PIL 로 실제 글자폭을 재어 줄 수를 세고, LibreOffice 의 넉넉한 행간을 반영한다.
from PIL import ImageFont  # noqa: E402

_LEAD = 1.02          # 실제 렌더로 보정한 계수 (line_spacing 과 곱해 실측 행간에 맞춘다)
_PAD_IN = 0.06        # 텍스트 상자 안쪽 여백
_fcache = {}


def _font(face, size_pt):
    key = (face, round(size_pt * 4))
    if key not in _fcache:
        fam = {T.DISPLAY: "KoPub Batang Bold.ttf",
               T.DISPLAY_MED: "KoPub Batang Medium.ttf",
               T.DISPLAY_LIGHT: "KoPub Batang Light.ttf",
               T.TEXT: "Pretendard-Regular.otf"}.get(face, "Pretendard-Regular.otf")
        _fcache[key] = ImageFont.truetype(str(fonts.find(fam)), int(size_pt * 4))
    return _fcache[key]


def measure(runs, width_in, size_pt, line_spacing=1.44, face=T.TEXT):
    """runs: 문자열 또는 [(문자열, opt), ...]. 필요한 높이(inch)를 돌려준다."""
    if isinstance(runs, str):
        runs = [(runs, {})]
    avail_px = (width_in - 0.10) * 72 * 4          # 4배 확대해 잰다
    lines, cur = 1, 0.0
    for txt, opt in [(r, {}) if isinstance(r, str) else r for r in runs]:
        f = _font(opt.get("face", face), opt.get("size", size_pt))
        for ch in txt:
            if ch == chr(10):
                lines += 1; cur = 0.0; continue
            w = f.getlength(ch)
            if cur + w > avail_px:
                lines += 1; cur = 0.0
            cur += w
    return lines * size_pt * line_spacing * _LEAD / 72.0 + _PAD_IN


# ══════════════════════════════════════════════════════════════════════════
# 1 ― 긴 문서형. 비대칭 2단. 왼쪽은 제도, 오른쪽은 우리 위치(신규성 선제 답변)
# ══════════════════════════════════════════════════════════════════════════
s = page(1)
text(s, "건설현장 AI CCTV를 어디에 달아야 하는지,\n정해 놓은 기준이 없다",
     T.ML, T.MT, T.CW - 0.9, 1.45, T.DISP, T.INK, face=T.DISPLAY, line=1.2)
ly1 = double_rule(s, T.ML, 2.16, T.CW) + 0.26

LEAD1 = ("LH 건설현장에는 AI CCTV가 달려 있다. 카메라가 찍은 화면을 인공지능이 보고, "
         "안전모를 안 쓴 사람을 찾아 관리자에게 알린다. 문제는 그 카메라를 어디에 다느냐에 "
         "따라 잘 찾기도 하고 못 찾기도 한다는 것이다.")
h_lead = measure(LEAD1, T.CW - 1.0, T.TXT + 1, 1.46)
text(s, LEAD1, T.ML, ly1, T.CW - 1.0, h_lead, T.TXT + 1, T.BODY, line=1.46)
y0 = ly1 + h_lead + 0.22

LW, GAP = 5.72, 0.62
RX = T.ML + LW + GAP
RW = T.CW - LW - GAP

text(s, "지금 기준이 정한 것과 안 정한 것", T.ML, y0, LW, 0.3, T.SUB, T.INK, face=T.DISPLAY)
text(s, "이 제안이 새로 하는 것", RX, y0, RW, 0.3, T.SUB, T.INK, face=T.DISPLAY)
yy = y0 + 0.40

LEFT = [
    ("성능은 인증서로 끝난다  ",
     "인증 시험은 미리 모아둔 영상에서 90% 이상 찾아내면 통과다. 정해진 영상으로만 "
     "재니까 카메라를 어디에 달아도 같은 점수가 나온다."),
    ("장비 사양은 있다  ",
     "AI CCTV 는 200만 화소 이상, 방수·방진 IP56 이상을 권장한다."),
    ("설치 조건은 없다  ",
     "국토교통부·국토안전관리원 가이드라인과 인증 기준 어디에도 카메라를 얼마나 멀리, "
     "어떤 각도로 달라는 숫자가 없다."),
]
RIGHT = [
    ("이미 있는 것  ",
     "확률로 카메라 자리를 정하는 방법은 이미 있다. 멀수록·비스듬할수록 못 찾는다는 계산, "
     "여러 대를 합치는 계산까지 모두 기존 연구의 것이다."),
    ("우리가 채우는 것  ",
     "그 확률을 재보지 않고 식으로 가정했다는 점이다. 직접 찍어서 재고, 가려짐을 0~100% "
     "사이 값으로 다루고, 심사할 때 계산서를 같이 보게 한다."),
    ("아직 확인 중  ",
     "앞선 연구 3건을 확인 중이다. 1건은 '직접 잰다'를 새롭지 않게 만들 수 있다."),
]

for col, wid, items in ((T.ML, LW, LEFT), (RX, RW, RIGHT)):
    cy = yy
    for i, (lab, txt) in enumerate(items):
        h = measure([(lab, {}), (txt, {})], wid, T.TXT, 1.44)
        if i:
            rule(s, col, cy - 0.10, wid)
        text(s, [[(lab, {"bold": True, "color": T.ACCENT}), (txt, {})]],
             col, cy, wid, h, T.TXT, T.BODY, line=1.44)
        cy += h + 0.22

note(s, "성능 기준은 KISA 지능형 CCTV 성능 시험인증 · 장비 사양은 국토교통부·국토안전관리원 "
        "「스마트 안전장비 가이드라인」(2026.5 개정) · 어느 문서에 무엇이 적혀 있는지는 원문 "
        "다시 확인 중")

# ══════════════════════════════════════════════════════════════════════════
# 2 ― 그림 지배형. 규정하는 축과 지배하는 축이 다르다는 것이 요점
# ══════════════════════════════════════════════════════════════════════════
s = page(2)
text(s, "지금 기준이 정한 것이 셋 중 영향이 가장 작았다",
     T.ML, T.MT, T.CW, 0.7, 32, T.INK, face=T.DISPLAY, line=1.18)
text(s, "셋은 단위가 달라 그대로 견줄 수 없다. 검출률이 90%까지 떨어지는 지점으로 "
        "맞춰 비교했다.",
     T.ML, 1.18, T.CW - 1.2, 0.32, T.TXT, T.MUTED, line=1.4)
rule(s, T.ML, 1.58, T.CW, weight=0.9)

fb = figure(s, "fig_curves.png", T.ML - 0.06, 1.72, 11.4)

STATS = [
    ("카메라와의 거리", f"{D_90_M:.0f}", " m",
     f"40 m에서는 아직 {PCT_AT_40M:.0f}%다", False),
    ("내려다보는 각도", f"{THETA_90:.0f}", "°",
     f"45°에서는 {PCT_AT_45DEG:.0f}%로 거의 그대로다", False),
    ("가려진 정도", f"{OCC_90:.1f}", "%",
     f"실제로 재본 15%에서는 이미 {PCT_AT_OCC15:.0f}%다", True),
]
sy = fb + 0.22
colw = (T.CW - 0.9) / 3
for i, (lab, num, unit, sub, hot) in enumerate(STATS):
    cx = T.ML + i * (colw + 0.45)
    if i:
        vrule(s, cx - 0.24, sy + 0.04, 1.02)
    col = T.ACCENT if hot else T.INK
    text(s, lab, cx, sy, colw, 0.28, T.LABEL, T.ACCENT if hot else T.MUTED)
    text(s, [[(num, {"face": T.DISPLAY, "size": T.STAT, "color": col}),
              (unit, {"face": T.DISPLAY_MED, "size": 17, "color": col})]],
         cx, sy + 0.28, colw, 0.56, T.STAT, col)
    text(s, sub, cx, sy + 0.84, colw, 0.3, T.FOLIO + 1, T.MUTED)

my = sy + 1.24
rule(s, T.ML, my, T.CW)
text(s, [[("어느 게 나쁘냐가 아니라 얼마나 나쁘냐  ", {"bold": True, "color": T.INK}),
          (f"가려지면 잘 못 찾는다는 건 다 아는 이야기다. 여기서 한 일은 그게 얼마나 "
           f"나쁜지를 {cp['n_conditions']}가지 조건에서 직접 재어, 배치 계산에 넣을 수 있는 "
           f"식으로 만든 것이다. 기존 연구는 가림을 막혔다·안 막혔다 둘 중 하나로만 다룬다. "
           f"위 세 값은 식으로 계산한 지점이고, 실제로 재본 가림은 15%부터다.", {})]],
     T.ML, my + 0.15, T.CW - 0.3, 0.66, T.FOLIO + 1, T.BODY, line=1.48)
note(s, f"안전모 사진 {cp['n_images']}장 · 세로축은 아무 조건도 안 걸었을 때({BASE_PCT:.1f}%)를 "
        f"1로 둔 비율 · 세 조건을 곱해 쓰는 식, 설명력 {cp['r2_full_grid']:.3f} · 검출기는 "
        f"안전모 사진으로 추가 학습시킨 YOLOv8n · 거리는 4K 해상도, 화각 90° 기준")

# ══════════════════════════════════════════════════════════════════════════
# 3 ― 문장 선언형. 평균이 붙어 있다는 것을 먼저 인정하고 들어간다
# ══════════════════════════════════════════════════════════════════════════
s = page(3)
text(s, "가정한 값으로 설계해도 평균은 비슷하다.\n차이는 사각지대에서 난다",
     T.ML, T.MT, T.CW - 1.2, 1.4, 33, T.INK, face=T.DISPLAY, line=1.2)
ly = double_rule(s, T.ML, 2.02, T.CW) + 0.22

text(s, [[("카메라 8대와 후보 자리는 그대로 두고 배치만 세 번 다르게 짰다. 셋 다 직접 잰 "
           "값으로 다시 채점했다. 사각지대는 ", {}),
          (f"{N_VOXEL:,}칸 중 {GEO['fail_voxel_count']} → {ASM['fail_voxel_count']} → "
           f"{EMP['fail_voxel_count']}칸", {"face": T.DISPLAY, "color": T.ACCENT}),
          ("으로 갈렸다.", {})]],
     T.ML + 0.9, ly, T.CW - 1.8, 0.92, T.LEAD - 4, T.BODY,
     face=T.DISPLAY_LIGHT, line=1.42)

fy = figure(s, "fig_three.png", (T.W - 10.7) / 2, ly + 1.00, 10.7)

by = fy + 0.16
rule(s, T.ML, by, T.CW)
CW2 = (T.CW - 0.8) / 2
text(s, [[("왜 갈리나  ", {"bold": True, "color": T.INK}),
          (f"가정한 값은 가림을 막혔다·안 막혔다 둘 중 하나로 본다. 조금 가려진 곳을 "
           f"괜찮다고 넘기는데, 그런 곳이 곧 사각지대다. 평균이 아니라 나쁜 쪽에서 갈린다"
           f"({abs(D_EA):.2f}%p 대 {ASM['fail_voxel_count'] - EMP['fail_voxel_count']}칸).",
           {})]],
     T.ML, by + 0.16, CW2, 1.1, T.FOLIO + 1, T.BODY, line=1.48)
text(s, [[("기준을 바꿔도  ", {"bold": True, "color": T.INK}),
          (f"사각지대로 치는 선을 {THR_SWEEP[0]['threshold']}·{THR_SWEEP[1]['threshold']}·"
           f"{THR_SWEEP[2]['threshold']}으로 바꿔도 제안 쪽이 늘 적다. 줄어드는 양은 "
           f"각각 {THR_SWEEP[0]['reduction']}·{THR_SWEEP[1]['reduction']}·"
           f"{THR_SWEEP[2]['reduction']}칸이며, {THR_SWEEP[2]['threshold']}에서는 차이가 "
           f"좁아진다. 가상 현장이라 숫자 자체가 아니라 세 배치의 순서를 주장한다.", {})]],
     T.ML + CW2 + 0.8, by + 0.16, CW2, 1.1, T.FOLIO + 1, T.BODY, line=1.48)
note(s, f"현장 {se['site']['width_m']}×{se['site']['depth_m']} m · 카메라 자리 후보 "
        f"{len(se['cameras'])}곳에서 {se['camera_budget']}대 · 비교 대상인 기존 방식은 국제표준 "
        f"네 등급을 다 돌려보고 기존 방식에 가장 유리한 등급으로 정했다 · 가정한 값은 국제표준의 "
        f"두 지점(125·250 PPM)으로 고정 · 한 대씩 고르는 방식이지만 최적해의 63% 이상이 "
        f"보장된다")

# ── 게이트 ─────────────────────────────────────────────────────────────────
def gate_glyphs():
    texts = [sh.text_frame.text for sl in prs.slides for sh in sl.shapes
             if sh.has_text_frame]
    fonts.assert_glyphs(texts, where="슬라이드")


def gate_font_slots():
    """모든 런에 latin·ea·cs 가 채워졌는지. ea 가 비면 한글만 테마 폰트로 나온다."""
    allowed = {T.DISPLAY, T.DISPLAY_MED, T.DISPLAY_LIGHT, T.TEXT}
    bad = 0
    for sl in prs.slides:
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    rPr = r._r.find(qn("a:rPr"))
                    if rPr is None:
                        bad += 1
                        continue
                    faces = {rPr.find(qn(t)).get("typeface")
                             if rPr.find(qn(t)) is not None else None
                             for t in ("a:latin", "a:ea", "a:cs")}
                    if len(faces) != 1 or faces.pop() not in allowed:
                        bad += 1
    if bad:
        raise SystemExit(f"폰트 슬롯이 비었거나 어긋난 런 {bad}개")


def gate_banned_words():
    """ADDENDUM-01 §4 금지 표현이 산출물에 새어 들어가는 것을 막는다."""
    banned = ["최초", "novel", "unprecedented", "기존에 없던", "낭비로 본다"]
    hits = []
    for i, sl in enumerate(prs.slides, 1):
        body = " ".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)
        hits += [f"S{i}:{w}" for w in banned if w in body]
    if hits:
        raise SystemExit("금지 표현(ADDENDUM-01 §4): " + " · ".join(hits))


gate_glyphs()
gate_font_slots()
gate_banned_words()
prs.save(str(OUT))
print(f"saved: {OUT.name} | slides: {len(prs.slides._sldIdLst)} | "
      f"display: {T.DISPLAY} | text: {T.TEXT}")
